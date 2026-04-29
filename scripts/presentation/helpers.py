"""
Helper functions for generating PowerPoint slide elements.
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from scripts.presentation.config import NAVY, BLUE, STEEL, WHITE, LIGHT, GRAY, DARK, ACCENT, GREEN, ORANGE, RED


def blank_slide(prs):
    """Add a blank slide to the presentation."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def bg(slide, color=NAVY):
    """Set the background color of a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=None, line=None):
    """Add a rectangular shape to the slide."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a textbox to the slide."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txb


def add_multiline(slide, lines, l, t, w, h, size=16, color=WHITE,
                  line_spacing=1.2, bold_first=False):
    """Add a multiline textbox. 'lines' is a list of str or (str, bool) for bold control."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, tuple):
            text, is_bold = line
        else:
            text, is_bold = line, (i == 0 and bold_first)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = color
        run.font.name = "Arial"
    return txb


def section_header(slide, title, subtitle="", accent=ACCENT):
    """Left accent bar + title on dark slide."""
    add_rect(slide, 0, 0, 0.5, 7.5, fill=accent)
    add_text(slide, title,  0.8, 2.8, 11.5, 1.0, size=36, bold=True,  color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.8, 3.9, 11.5, 0.8, size=20, color=ACCENT)


def kpi_box(slide, value, label, l, t, w=2.5, h=1.4, vcolor=ACCENT, bg_color=NAVY):
    """Draw a KPI display box."""
    add_rect(slide, l, t, w, h, fill=bg_color, line=ACCENT)
    add_text(slide, value, l+0.1, t+0.05, w-0.2, 0.75,
             size=28, bold=True, color=vcolor, align=PP_ALIGN.CENTER)
    add_text(slide, label, l+0.1, t+0.8, w-0.2, 0.55,
             size=12, color=WHITE, align=PP_ALIGN.CENTER)


def table_box(slide, headers, rows, l, t, w, h,
              hdr_fill=NAVY, hdr_color=WHITE, row_fill=DARK, row_color=WHITE,
              alt_fill=None, font_size=12):
    """Draw a table layout."""
    col_n = len(headers)
    col_w = w / col_n
    row_h = h / (len(rows) + 1)

    # header row
    for ci, hdr in enumerate(headers):
        add_rect(slide, l + ci*col_w, t, col_w, row_h, fill=hdr_fill)
        add_text(slide, hdr, l + ci*col_w + 0.05, t + 0.02,
                 col_w - 0.1, row_h - 0.04, size=font_size, bold=True, color=hdr_color,
                 align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        fill = (alt_fill if (alt_fill and ri % 2 == 1) else row_fill)
        for ci, cell in enumerate(row):
            add_rect(slide, l + ci*col_w, t + (ri+1)*row_h, col_w, row_h, fill=fill)
            add_text(slide, str(cell),
                     l + ci*col_w + 0.05, t + (ri+1)*row_h + 0.02,
                     col_w - 0.1, row_h - 0.04, size=font_size-1, color=row_color,
                     align=PP_ALIGN.CENTER)


def flow_box(slide, steps, l, t, box_w=1.7, box_h=0.7, gap=0.2, color=BLUE):
    """Draw horizontal flow boxes with arrows."""
    x = l
    for i, (num, label) in enumerate(steps):
        add_rect(slide, x, t, box_w, box_h, fill=color)
        add_text(slide, f"{num}", x+0.05, t+0.02, 0.3, 0.3, size=11, bold=True, color=ACCENT)
        add_text(slide, label, x+0.05, t+0.3, box_w-0.1, box_h-0.3, size=11, color=WHITE)
        if i < len(steps)-1:
            add_text(slide, "→", x+box_w+0.02, t+0.15, gap+0.1, 0.4,
                     size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        x += box_w + gap + 0.1
