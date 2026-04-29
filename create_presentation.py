"""
Generate Thai NLP-to-SQL Agent PowerPoint presentation from report content.
"""
from pptx import Presentation
import os

from scripts.presentation.config import W, H
from scripts.presentation.slides import build_all_slides

def main():
    # Initialize presentation
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Build slides
    build_all_slides(prs)

    # Save presentation
    base_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(base_dir, "presentation.pptx")
    
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides generated: {len(prs.slides)}")

if __name__ == "__main__":
    main()
